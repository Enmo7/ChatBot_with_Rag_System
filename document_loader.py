import os
import glob
import fitz
import re
import gc
import psutil
import tempfile
import time
import hashlib
from typing import Iterator
from PIL import Image
from docx import Document as DocxDocument
from pptx import Presentation
from langchain_community.document_loaders import TextLoader
from langchain_core.documents import Document
from rapidocr_onnxruntime import RapidOCR
from langchain_text_splitters import RecursiveCharacterTextSplitter
from metadata_store import MetadataStore

class DocumentLoader:
    def __init__(self, directory_path="./documents"):
        self.directory_path = os.path.normpath(directory_path)
        self.metadata_store = MetadataStore()
        
        # ✅ CRITICAL FIX: Force CPU mode to avoid GPU permission issues
        use_gpu = False  # Set to False to use CPU (avoids WinError 5 permission issues)
        
        # ✅ CRITICAL FIX: Handle model download issues causing 'KeyError: model_path'
        try:
            # Initialize with explicit CPU parameters AND force model download
            self.ocr_engine = RapidOCR(
                det_use_cuda=use_gpu,
                cls_use_cuda=use_gpu,
                rec_use_cuda=use_gpu,
                # Force model download on initialization
                det_model_path=None,
                cls_model_path=None,
                rec_model_path=None
            )
            device_msg = "GPU (CUDA)" if use_gpu else "CPU"
            print(f"✅ OCR engine initialized on {device_msg} using RapidOCR")
        except Exception as e:
            self.ocr_engine = None
            error_msg = str(e)
            # ✅ Handle 'model_path' KeyError specifically
            if "model_path" in error_msg or "KeyError" in error_msg:
                print("⚠️ OCR initialization failed: Missing model files")
                print("💡 FIX: Models will auto-download on first use. Restart the application.")
                print("   If issue persists:")
                print("   1. pip uninstall rapidocr-onnxruntime -y")
                print("   2. pip install rapidocr-onnxruntime --no-cache-dir")
                print("   3. Restart your Python environment")
            else:
                print(f"⚠️ OCR engine initialization failed - falling back to text-only extraction. Error: {type(e).__name__}: {e}")
                print("💡 Troubleshooting:")
                print("   1. Install CPU runtime: pip install onnxruntime")
                print("   2. Force CPU mode: Already enabled in code (use_gpu=False)")
                print("   3. Restart application after package installation")
        
        self.text_splitter = RecursiveCharacterTextSplitter(chunk_size=2000, chunk_overlap=300)
        
        # ✅ CRITICAL FIX: Removed space before REQ pattern - was causing regex to fail
        self.link_pattern = re.compile(r'\b(REQ|TC|SRS|TEST|BUG|SPEC|DOC|FTR)-?[A-Z0-9]+\b', re.IGNORECASE)

    def check_memory(self):
        """Pauses execution if memory is high. Returns True if safe to proceed."""
        if psutil.virtual_memory().percent > 90:
            print("⚠️ High Memory! Pausing for garbage collection...")
            gc.collect()
            time.sleep(5)
            if psutil.virtual_memory().percent > 95:
                print("❌ Memory Critical! Skipping current batch.")
                return False
        return True

    def hash_chunk(self, text):
        """Generates MD5 hash for chunk integrity verification"""
        return hashlib.md5(text.encode('utf-8')).hexdigest()

    def extract_links_with_full_info(self, text):
        raw_links = list(set(self.link_pattern.findall(text)))
        results = []
        for link in raw_links:
            idx = text.find(link)
            if idx == -1: 
                continue
            start, end = max(0, idx - 50), min(len(text), idx + len(link) + 50)
            context = text[start:end]  # ✅ FIXED: 'contex t' → 'context'
            results.append({
                "id": link, 
                "context": context, 
                "valid_context": len(context.split()) > 5
            })
        return results

    def run_ocr_safely(self, image_source):
        if not self.ocr_engine: 
            return ""
        temp_file = None
        try:
            temp_file = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
            temp_file.close()
            if isinstance(image_source, bytes):
                with open(temp_file.name, 'wb') as f: 
                    f.write(image_source)
            else: 
                image_source.save(temp_file.name)
            
            with open(temp_file.name, 'rb') as f: 
                res, _ = self.ocr_engine(f.read())
            return "\n".join([line[1] for line in res]) if res else ""
        except Exception as e:
            print(f"OCR error: {e}")
            return ""
        finally:
            if temp_file and os.path.exists(temp_file.name): 
                os.unlink(temp_file.name)

    def _process_text_content(self, text, file_hash, file_path, **metadata):
        if not text.strip():  # ✅ FIXED: 'n ot' → 'not'
            return
        
        links_info = self.extract_links_with_full_info(text)
        valid_links = []
        for link in links_info:
            self.metadata_store.log_detected_link(file_hash, link['id'], link['context'], 1.0)
            if link['valid_context']: 
                valid_links.append(link['id'])

        doc = Document(page_content=text, metadata={"source": file_path, "links": ", ".join(valid_links), **metadata})  # ✅ FIXED: 'm etadata' → 'metadata'
        chunks = self.text_splitter.split_documents([doc])
        for chunk in chunks:
            chunk.metadata['chunk_hash'] = self.hash_chunk(chunk.page_content)  # ✅ FIXED: 'page_ content' → 'page_content'
            yield chunk

    def process_file_generator(self) -> Iterator[Document]:
        if not os.path.exists(self.directory_path): 
            return
        
        for file_path in glob.glob(os.path.join(self.directory_path, "*.*")):
            try:
                if not self.check_memory(): 
                    continue

                file_hash, is_new = self.metadata_store.register_document(file_path, os.path.basename(file_path))  # ✅ FIXED: 'register_d ocument' → 'register_document'
                if not file_hash: 
                    continue
                
                ext = file_path.rsplit('.', 1)[-1].lower()
                
                if ext == 'pdf':
                    doc = fitz.open(file_path)
                    for i, page in enumerate(doc):
                        if i % 20 == 0: 
                            self.check_memory()
                        text = page.get_text()
                        if len(text) < 50:  # OCR fallback for scanned pages
                            pix = page.get_pixmap(matrix=fitz.Matrix(1.0, 1.0))
                            text = self.run_ocr_safely(pix.tobytes("png"))
                        yield from self._process_text_content(text, file_hash, file_path, page=i+1, type="pdf")
                    doc.close()

                elif ext == 'pptx':
                    prs = Presentation(file_path)
                    for i, slide in enumerate(prs.slides):
                        text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
                        yield from self._process_text_content(text, file_hash, file_path, page=i+1, type="slide")

                elif ext in ['png', 'jpg', 'jpeg']:
                    with Image.open(file_path) as img:
                        text = self.run_ocr_safely(img)
                        yield from self._process_text_content(text, file_hash, file_path, type="image")

                elif ext == 'docx':
                    doc = DocxDocument(file_path)
                    text = "\n".join([p.text for p in doc.paragraphs])
                    yield from self._process_text_content(text, file_hash, file_path, type="docx")
                    
                elif ext == 'txt':
                    loader = TextLoader(file_path, encoding='utf-8')
                    for d in loader.load():
                        yield from self._process_text_content(d.page_content, file_hash, file_path, type="text")  # ✅ FIXED: 'yiel d' → 'yield'

            except Exception as e:
                print(f"Error processing {file_path}: {e}")
                import traceback
                traceback.print_exc()